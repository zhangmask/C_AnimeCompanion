import os
import shutil
import tempfile
import uuid

import pytest
from build_test_helpers import (
    assert_content_no_html_tags,
    assert_resource_indexed,
    assert_root_uri_valid,
    assert_source_format,
)


def _create_html_file():
    random_id = str(uuid.uuid4())[:8]
    unique_keyword = f"html_keyword_{random_id}"
    temp_dir = tempfile.mkdtemp()
    html_path = os.path.join(temp_dir, f"test_{random_id}.html")

    content = (
        f"<html><head><title>HTML Test {random_id}</title></head>"
        f"<body><h1>HTML测试标题 {random_id}</h1>"
        f"<p>包含唯一关键词：{unique_keyword}</p>"
        f"<p>用于验证HTML文件构建产物。</p>"
        f"</body></html>"
    )
    with open(html_path, "w", encoding="utf-8") as f:
        f.write(content)

    return html_path, temp_dir, unique_keyword


def _create_pptx_file():
    try:
        from pptx import Presentation
    except ImportError:
        pytest.skip("python-pptx 未安装，跳过 PPTX 构建测试")

    random_id = str(uuid.uuid4())[:8]
    unique_keyword = f"pptx_keyword_{random_id}"
    temp_dir = tempfile.mkdtemp()
    pptx_path = os.path.join(temp_dir, f"test_{random_id}.pptx")

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = f"PPT测试标题 {random_id}"
    body = slide.placeholders[1]
    body.text = f"包含唯一关键词：{unique_keyword}\n用于验证PPT文件构建产物。"
    prs.save(pptx_path)

    return pptx_path, temp_dir, unique_keyword


class TestBuildDocumentResources:
    """TC-B04, B07 文档类资源构建测试（快速用例，≤20s）"""

    def test_build_html_file(self, api_client):
        """TC-B04 HTML文件构建：验证 .html 文件添加后 source_format=html 且标签被剥离"""
        html_path, temp_dir, unique_keyword = _create_html_file()
        try:
            response = api_client.add_resource(path=html_path, wait=True)
            assert response.status_code == 200

            data = response.json()
            assert data.get("status") == "ok"

            result = data.get("result", {})
            root_uri = result.get("root_uri")
            assert_root_uri_valid(root_uri)

            stat_resp = api_client.fs_stat(root_uri)
            assert stat_resp.status_code == 200

            assert_source_format(api_client, root_uri, ["html", "markdown"])

            assert_content_no_html_tags(api_client, root_uri)

            assert_resource_indexed(api_client, root_uri, unique_keyword)

            print(f"✓ TC-B04 HTML文件构建通过, root_uri: {root_uri}")
        finally:
            shutil.rmtree(temp_dir, ignore_errors=True)

    def test_build_pptx_file(self, api_client):
        """TC-B07 PowerPoint构建：验证 .pptx 文件添加后 source_format=pptx 且内容可检索"""
        pptx_path, temp_dir, unique_keyword = _create_pptx_file()
        try:
            response = api_client.add_resource(path=pptx_path, wait=True)
            assert response.status_code == 200

            data = response.json()
            assert data.get("status") == "ok"

            result = data.get("result", {})
            root_uri = result.get("root_uri")
            assert_root_uri_valid(root_uri)

            stat_resp = api_client.fs_stat(root_uri)
            assert stat_resp.status_code == 200

            assert_source_format(api_client, root_uri, ["pptx", "markdown"])

            assert_resource_indexed(api_client, root_uri, unique_keyword)

            print(f"✓ TC-B07 PowerPoint构建通过, root_uri: {root_uri}")
        finally:
            shutil.rmtree(temp_dir, ignore_errors=True)
