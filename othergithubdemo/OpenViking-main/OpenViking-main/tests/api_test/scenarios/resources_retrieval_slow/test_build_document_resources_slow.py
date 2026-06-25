import os
import shutil
import tempfile
import uuid

import pytest
from build_test_helpers import (
    assert_resource_indexed,
    assert_root_uri_valid,
    assert_source_format,
)


def _create_pdf_file():
    try:
        from fpdf import FPDF
    except ImportError:
        pytest.skip("fpdf 未安装，跳过 PDF 构建测试")

    random_id = str(uuid.uuid4())[:8]
    unique_keyword = f"pdf_keyword_{random_id}"
    temp_dir = tempfile.mkdtemp()
    pdf_path = os.path.join(temp_dir, f"test_{random_id}.pdf")

    pdf = FPDF()
    pdf.add_page()
    pdf.set_font("Helvetica", size=12)
    pdf.cell(200, 10, text=f"PDF Test Document {random_id}", new_x="LMARGIN", new_y="NEXT")
    pdf.cell(200, 10, text=f"Unique keyword: {unique_keyword}", new_x="LMARGIN", new_y="NEXT")
    pdf.cell(
        200, 10, text="This is a test PDF file for build validation.", new_x="LMARGIN", new_y="NEXT"
    )
    pdf.output(pdf_path)

    return pdf_path, temp_dir, unique_keyword


def _create_docx_file():
    try:
        from docx import Document
    except ImportError:
        pytest.skip("python-docx 未安装，跳过 DOCX 构建测试")

    random_id = str(uuid.uuid4())[:8]
    unique_keyword = f"docx_keyword_{random_id}"
    temp_dir = tempfile.mkdtemp()
    docx_path = os.path.join(temp_dir, f"test_{random_id}.docx")

    doc = Document()
    doc.add_heading(f"Word测试标题 {random_id}", level=1)
    doc.add_paragraph(f"包含唯一关键词：{unique_keyword}")
    doc.add_paragraph("用于验证Word文档构建产物。")
    doc.save(docx_path)

    return docx_path, temp_dir, unique_keyword


def _create_xlsx_file():
    try:
        from openpyxl import Workbook
    except ImportError:
        pytest.skip("openpyxl 未安装，跳过 XLSX 构建测试")

    random_id = str(uuid.uuid4())[:8]
    unique_keyword = f"xlsx_keyword_{random_id}"
    temp_dir = tempfile.mkdtemp()
    xlsx_path = os.path.join(temp_dir, f"test_{random_id}.xlsx")

    wb = Workbook()
    ws = wb.active
    ws.title = "TestSheet"
    ws.append(["Column A", "Column B", "Column C"])
    ws.append([f"数据1 {random_id}", unique_keyword, "数据3"])
    ws.append(["数据4", "数据5", "数据6"])
    wb.save(xlsx_path)

    return xlsx_path, temp_dir, unique_keyword


class TestBuildDocumentResourcesSlow:
    """TC-B03~B09 文档类资源构建测试"""

    def test_build_pdf_file(self, api_client):
        """TC-B03 PDF文件构建：验证 .pdf 文件添加后 source_format=pdf 且内容可检索"""
        pdf_path, temp_dir, unique_keyword = _create_pdf_file()
        try:
            response = api_client.add_resource(path=pdf_path, wait=True)
            assert response.status_code == 200

            data = response.json()
            assert data.get("status") == "ok"

            result = data.get("result", {})
            root_uri = result.get("root_uri")
            assert_root_uri_valid(root_uri)

            stat_resp = api_client.fs_stat(root_uri)
            assert stat_resp.status_code == 200

            assert_source_format(api_client, root_uri, ["pdf", "markdown"])

            assert_resource_indexed(api_client, root_uri, unique_keyword)

            print(f"✓ TC-B03 PDF文件构建通过, root_uri: {root_uri}")
        finally:
            shutil.rmtree(temp_dir, ignore_errors=True)

    def test_build_docx_file(self, api_client):
        """TC-B05 Word文档构建：验证 .docx 文件添加后 source_format=docx 且内容可检索"""
        docx_path, temp_dir, unique_keyword = _create_docx_file()
        try:
            response = api_client.add_resource(path=docx_path, wait=True)
            assert response.status_code == 200

            data = response.json()
            assert data.get("status") == "ok"

            result = data.get("result", {})
            root_uri = result.get("root_uri")
            assert_root_uri_valid(root_uri)

            stat_resp = api_client.fs_stat(root_uri)
            assert stat_resp.status_code == 200

            assert_source_format(api_client, root_uri, ["docx", "markdown"])

            assert_resource_indexed(api_client, root_uri, unique_keyword)

            print(f"✓ TC-B05 Word文档构建通过, root_uri: {root_uri}")
        finally:
            shutil.rmtree(temp_dir, ignore_errors=True)

    def test_build_xlsx_file(self, api_client):
        """TC-B08 Excel构建：验证 .xlsx 文件添加后 source_format=xlsx 且表格数据可检索"""
        xlsx_path, temp_dir, unique_keyword = _create_xlsx_file()
        try:
            response = api_client.add_resource(path=xlsx_path, wait=True)
            assert response.status_code == 200

            data = response.json()
            assert data.get("status") == "ok"

            result = data.get("result", {})
            root_uri = result.get("root_uri")
            assert_root_uri_valid(root_uri)

            stat_resp = api_client.fs_stat(root_uri)
            assert stat_resp.status_code == 200

            assert_source_format(api_client, root_uri, ["xlsx", "markdown"])

            assert_resource_indexed(api_client, root_uri, unique_keyword)

            print(f"✓ TC-B08 Excel构建通过, root_uri: {root_uri}")
        finally:
            shutil.rmtree(temp_dir, ignore_errors=True)

    def test_build_html_file(self, api_client):
        """TC-B04 HTML文件构建：验证 .html 文件添加后 source_format=html 且标签被剥离"""
        from build_test_helpers import assert_content_no_html_tags

        random_id = str(uuid.uuid4())[:8]
        unique_keyword = f"html_keyword_{random_id}"
        temp_dir = tempfile.mkdtemp()
        html_path = os.path.join(temp_dir, f"test_{random_id}.html")

        content = (
            f"<html><head><title>HTML Test {random_id}</title></head>"
            f"<body><h1>HTML测试标题 {random_id}</h1>"
            f"<p>包含唯一关键词：{unique_keyword}</p>"
            f"<p>用于验证HTML文件构建产物。</p>"
            f"</body></html>"
        )
        with open(html_path, "w", encoding="utf-8") as f:
            f.write(content)

        try:
            response = api_client.add_resource(path=html_path, wait=True)
            assert response.status_code == 200

            data = response.json()
            assert data.get("status") == "ok"

            result = data.get("result", {})
            root_uri = result.get("root_uri")
            assert_root_uri_valid(root_uri)

            stat_resp = api_client.fs_stat(root_uri)
            assert stat_resp.status_code == 200

            assert_source_format(api_client, root_uri, ["html", "markdown"])

            assert_content_no_html_tags(api_client, root_uri)

            assert_resource_indexed(api_client, root_uri, unique_keyword)

            print(f"✓ TC-B04 HTML文件构建通过, root_uri: {root_uri}")
        finally:
            shutil.rmtree(temp_dir, ignore_errors=True)

    def test_build_legacy_doc_file(self, api_client):
        """TC-B06 旧版Word(.doc)构建：验证 .doc 文件添加后 source_format=doc 且内容可检索"""
        random_id = str(uuid.uuid4())[:8]
        unique_keyword = f"doc_keyword_{random_id}"
        temp_dir = tempfile.mkdtemp()
        doc_path = os.path.join(temp_dir, f"test_{random_id}.doc")

        try:
            from docx import Document

            doc = Document()
            doc.add_heading(f"旧版Word测试标题 {random_id}", level=1)
            doc.add_paragraph(f"包含唯一关键词：{unique_keyword}")
            doc.save(doc_path)
        except ImportError:
            pytest.skip("python-docx 未安装，跳过旧版Word构建测试")

        try:
            response = api_client.add_resource(path=doc_path, wait=True)
            assert response.status_code == 200

            data = response.json()
            assert data.get("status") == "ok"

            result = data.get("result", {})
            root_uri = result.get("root_uri")
            assert_root_uri_valid(root_uri)

            stat_resp = api_client.fs_stat(root_uri)
            assert stat_resp.status_code == 200

            assert_source_format(api_client, root_uri, ["doc", "docx", "markdown"])

            assert_resource_indexed(api_client, root_uri, unique_keyword)

            print(f"✓ TC-B06 旧版Word(.doc)构建通过, root_uri: {root_uri}")
        finally:
            shutil.rmtree(temp_dir, ignore_errors=True)

    def test_build_pptx_file(self, api_client):
        """TC-B07 PowerPoint构建：验证 .pptx 文件添加后 source_format=pptx 且内容可检索"""
        try:
            from pptx import Presentation
        except ImportError:
            pytest.skip("python-pptx 未安装，跳过 PPTX 构建测试")

        random_id = str(uuid.uuid4())[:8]
        unique_keyword = f"pptx_keyword_{random_id}"
        temp_dir = tempfile.mkdtemp()
        pptx_path = os.path.join(temp_dir, f"test_{random_id}.pptx")

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title = slide.shapes.title
        title.text = f"PPT测试标题 {random_id}"
        body = slide.placeholders[1]
        body.text = f"包含唯一关键词：{unique_keyword}\n用于验证PPT文件构建产物。"
        prs.save(pptx_path)

        try:
            response = api_client.add_resource(path=pptx_path, wait=True)
            assert response.status_code == 200

            data = response.json()
            assert data.get("status") == "ok"

            result = data.get("result", {})
            root_uri = result.get("root_uri")
            assert_root_uri_valid(root_uri)

            stat_resp = api_client.fs_stat(root_uri)
            assert stat_resp.status_code == 200

            assert_source_format(api_client, root_uri, ["pptx", "markdown"])

            assert_resource_indexed(api_client, root_uri, unique_keyword)

            print(f"✓ TC-B07 PowerPoint构建通过, root_uri: {root_uri}")
        finally:
            shutil.rmtree(temp_dir, ignore_errors=True)

    def test_build_epub_file(self, api_client):
        """TC-B09 EPUB构建：验证 .epub 文件添加后 source_format=epub 且内容可检索"""
        try:
            from ebooklib import epub
        except ImportError:
            pytest.skip("ebooklib 未安装，跳过 EPUB 构建测试")

        random_id = str(uuid.uuid4())[:8]
        unique_keyword = f"epub_keyword_{random_id}"
        temp_dir = tempfile.mkdtemp()
        epub_path = os.path.join(temp_dir, f"test_{random_id}.epub")

        book = epub.EpubBook()
        book.set_identifier(f"test-epub-{random_id}")
        book.set_title(f"EPUB测试 {random_id}")
        book.set_language("zh")

        chapter = epub.EpubHtml(title="Chapter 1", file_name="chap_01.xhtml", lang="zh")
        chapter.content = (
            f"<html><body><h1>EPUB测试章节 {random_id}</h1>"
            f"<p>包含唯一关键词：{unique_keyword}</p>"
            f"<p>用于验证EPUB文件构建产物。</p>"
            f"</body></html>"
        )
        book.add_item(chapter)
        book.toc = (chapter,)
        book.add_item(epub.EpubNcx())
        book.add_item(epub.EpubNav())
        epub.write_epub(epub_path, book, {})

        try:
            response = api_client.add_resource(path=epub_path, wait=True)
            assert response.status_code == 200

            data = response.json()
            assert data.get("status") == "ok"

            result = data.get("result", {})
            root_uri = result.get("root_uri")
            assert_root_uri_valid(root_uri)

            stat_resp = api_client.fs_stat(root_uri)
            assert stat_resp.status_code == 200

            assert_source_format(api_client, root_uri, ["epub", "markdown"])

            assert_resource_indexed(api_client, root_uri, unique_keyword)

            print(f"✓ TC-B09 EPUB构建通过, root_uri: {root_uri}")
        finally:
            shutil.rmtree(temp_dir, ignore_errors=True)
